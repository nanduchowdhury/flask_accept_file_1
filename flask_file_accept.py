
#########################
#
# Install packages in google-shell as follows:
#   /usr/bin/pip3.10 install PyPDF2
#########################

from flask import Flask, render_template, request, jsonify
from werkzeug.utils import secure_filename
import os

import PyPDF2

import pathlib
import textwrap

import base64
from io import BytesIO
import io
from PIL import Image, ImageDraw

from flask import Flask, request, send_file
from pptx import Presentation
from pptx.util import Inches  # Import Inches from pptx.util
from io import BytesIO
import pdfkit  # You'll need to install the `pdfkit` and `wkhtmltopdf

import tempfile  # For creating temporary files

from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

import google.generativeai as genai

from error_message_manager import ErrorManager
from gemini_access import GeminiAccess
from base_client_manager import BaseClientManager


class ScholarKM(Flask, BaseClientManager):
    def __init__(self):
        Flask.__init__(self, __name__)
        BaseClientManager.__init__(self)  # Initialize client management base class

        self.client_ip = 'client_ip'

        self.route('/')(self.index)
        self.route('/upload', methods=['POST'])(self.upload_file)
        self.route('/save_logs', methods=['POST'])(self.save_logs)
        self.route('/explain_region', methods=['POST'])(self.explain_region)
        self.route('/convert_ppt_to_pdf', methods=['POST'])(self.convert_ppt_to_pdf)

        # Constants
        self.BASE_FOLDER = '/home/nandu_chowdhury/kupamanduk/scholar/'
        self.SERVER_LOGS_FOLDER = '/home/nandu_chowdhury/kupamanduk/scholar/server_logs/'
        self.BASE_UPLOAD_FOLDER = 'uploads/'
        self.BASE_LOG_FOLDER = 'client_logs/'
        self.FIRST_TITLE_LEVEL = 0
        self.ALL_TITLES_LEVEL = -1

        # Key names in the shared-data dictionary
        self.CDATA_client_id = 'client_id'
        self.CDATA_client_ip = 'client_ip'
        self.CDATA_upload_folder = 'upload_folder'
        self.CDATA_log_folder = 'log_folder'
        self.CDATA_log_file = 'log_file'
        self.CDATA_main_content_file = 'main_content_file'
        self.CDATA_num_learn_points = 'num_learn_points'
        self.CDATA_first_response = 'first_response'

        self.error_manager = ErrorManager(self.client_ip, 'static/errors.txt', 
                    log_dir=self.SERVER_LOGS_FOLDER)

        self.gemini_access = GeminiAccess(self.error_manager)
        self.gemini_access.initialize()

    def extract_text_from_pdf(self, pdf_path):
        text = ""
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page_num in range(len(reader.pages)):
                page = reader.pages[page_num]
                text += page.extract_text()
        return text

    def establish_folders(self, uuid):

        client_ip = self.get_client_data(uuid, self.CDATA_client_ip)
        client_id = self.get_client_data(uuid, self.CDATA_client_id)
        
        client_folder = os.path.join(self.BASE_FOLDER, f'{client_ip}_{client_id}/')

        if not os.path.exists(client_folder):
            os.makedirs(client_folder, exist_ok=True)

        upload_folder = os.path.join(client_folder, self.BASE_UPLOAD_FOLDER)
        self.save_client_data(uuid, self.CDATA_upload_folder, upload_folder)

        if not os.path.exists(upload_folder):
            os.makedirs(upload_folder, exist_ok=True)
        
        log_folder = os.path.join(client_folder, self.BASE_LOG_FOLDER)
        self.save_client_data(uuid, self.CDATA_log_folder, log_folder)

        if not os.path.exists(log_folder):
            os.makedirs(log_folder, exist_ok=True)
        
        log_file_path = os.path.join(log_folder, f'client_log.txt')
        self.save_client_data(uuid, self.CDATA_log_file, log_file_path)

    def extract_file(self, uuid, data):

        file_name = data['fileName']

        upload_folder = self.get_client_data(uuid, self.CDATA_upload_folder)

        file_name = os.path.join(upload_folder, file_name)
        file_content = base64.b64decode(data['fileContent'])

        with open(file_name, 'wb') as f:
            f.write(file_content)
        return file_name

    def extract_image(self, uuid, data):

        upload_folder = self.get_client_data(uuid, self.CDATA_upload_folder)

        image_data_url = data['image']
        header, base64_image = image_data_url.split(',', 1)
        file_content = base64.b64decode(base64_image)
        image = Image.open(io.BytesIO(file_content))
        filename = os.path.join(upload_folder, 'captured_image.png')
        image.save(filename)
        return filename

    def extract_video(self, uuid, data):

        upload_folder = self.get_client_data(uuid, self.CDATA_upload_folder)

        video_data = data['video']
        video_binary = base64.b64decode(video_data)
        video_path = os.path.join(upload_folder, 'uploaded_video.webm')
        with open(video_path, 'wb') as video_file:
            video_file.write(video_binary)
        return video_path

    def index(self):
        return render_template('index.html')

    def receive_and_save_file(self, uuid, data):
        file_path = ''

        if 'fileContent' in data:
            file_path = self.extract_file(uuid, data)

            self.error_manager.show_message(2002, file_path)

        elif 'image' in data:
            file_path = self.extract_image(uuid, data)
            self.error_manager.show_message(2003, file_path)

        elif 'video' in data:
            file_path = self.extract_video(uuid, data)
            self.error_manager.show_message(2004, file_path)

        else:
            raise ValueError(self.error_manager.show_message(2001))

        return file_path

    def upload_file(self):
        
        try:
            data = request.json

            uuid = self.get_or_generate_uuid(data)

            learnLevel = data['additionalData']
            learnLevel = learnLevel['learnLevel']

            if learnLevel == self.ALL_TITLES_LEVEL:
            
                main_content_file = self.receive_and_save_file(uuid, data)
                self.save_client_data(uuid, self.CDATA_main_content_file, main_content_file)
                self.gemini_access.upload_file(uuid, main_content_file)

                self.gemini_access.check_content_student_related(uuid)

                if ( self.gemini_access.is_there_text_in_content(uuid) ):
                    self.error_manager.show_message(2012, "TEXT")
                    english_response = self.gemini_access.get_all_headers_of_text(uuid)
                else:
                    self.error_manager.show_message(2012, "PICTURE")
                    english_response = self.gemini_access.get_all_headers_of_picture(uuid)

                num_learn_points = len(english_response.splitlines())
                self.save_client_data(uuid, self.CDATA_num_learn_points, main_content_file)

                self.error_manager.show_message(2009, learnLevel, num_learn_points)

                self.save_client_data(uuid, self.CDATA_first_response, english_response)

                first_response = english_response.splitlines()

                hindi_response = self.gemini_access.convert_to_hindi(english_response)

                return jsonify({'numPoints': num_learn_points,
                                'firstResponse' : first_response,
                                'result1': english_response, 
                                'result2': hindi_response})
            else:

                first_response = self.get_client_data(uuid, self.CDATA_first_response)

                first_response = first_response.splitlines()

                english_response = self.gemini_access.get_header_summary(uuid, first_response[learnLevel])
                hindi_response = self.gemini_access.convert_to_hindi(english_response)
                self.error_manager.show_message(2010, learnLevel)

                return jsonify({'result1': english_response, 
                                'result2': hindi_response})

        except Exception as e:
            return jsonify({"error": str(e)}), 500


    def explain_region(self):
        try:
            data = request.get_json()
            
            uuid = self.get_or_generate_uuid(data)

            main_content_file = self.get_client_data(uuid, self.CDATA_main_content_file)
            explain_region_file = self.extract_image(data)

            english_response = self.gemini_access.explain_region(uuid, main_content_file, explain_region_file)
            self.error_manager.show_message(2013)

            hindi_response = self.gemini_access.convert_to_hindi(english_response)

            return jsonify({'result1': english_response, 
                            'result2': hindi_response})

        except Exception as e:
            return jsonify({"error": str(e)}), 500

    def get_or_generate_uuid(self, data, is_generate=False):
        client_uuid = data.get('client_uuid', '')
        print("UUID : ", client_uuid)
        if not client_uuid:
            if is_generate:
                client_uuid = self.set_client_id()  # Generate a new client ID
            else:
                raise ValueError(self.error_manager.show_message(2017))
        return client_uuid

    def save_logs(self):
        try:
            data = request.get_json()
            
            client_uuid = self.get_or_generate_uuid(data, True)
            client_id = data.get('clientId', 'unknown')
            self.save_client_data(client_uuid, self.CDATA_client_id, client_id)
            
            client_ip = request.remote_addr
            self.save_client_data(client_uuid, self.CDATA_client_ip, client_ip)

            self.error_manager.update_client_ip(client_ip)

            self.establish_folders(client_uuid)

            logs = data.get('logs', [])

            if not logs:
                return jsonify({"status": "no logs to save"}), 200
            
            log_file_path = self.get_client_data(client_uuid, self.CDATA_log_file)

            # Append the logs to the client's log file
            with open(log_file_path, 'a') as log_file:
                for log in logs:
                    log_file.write(f"{log}\n")


            return jsonify({"client_uuid": client_uuid}), 200

        except Exception as e:
            return jsonify({"error": str(e)}), 500



    def convert_slide_to_image(self, slide):
        # For simplicity, we'll create a blank image to represent the slide
        width, height = 1280, 720  # Default slide dimensions in pixels

        # Create a blank white image for the slide
        slide_image = Image.new('RGB', (width, height), color='white')
        draw = ImageDraw.Draw(slide_image)
        
        # Iterate over the shapes in the slide and draw rectangles and text
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            # Extract text from the shape
            text = shape.text
            
            # Convert EMUs to pixels (assuming 96 DPI)
            left = shape.left / Inches(1) * 96  # Convert Inches to pixels
            top = shape.top / Inches(1) * 96
            width = shape.width / Inches(1) * 96
            height = shape.height / Inches(1) * 96
            
            # Draw a rectangle around the shape
            draw.rectangle([left, top, left + width, top + height], outline='black', width=2)
            
            # Draw the text inside the rectangle
            draw.text((left + 10, top + 10), text, fill='black')
        
        return slide_image

    def convert_ppt_to_pdf(self):
        try:

            if 'file' not in request.files:
                raise ValueError(self.error_manager.show_message(2016))

            file = request.files['file']
            
            if file.filename.endswith('.pptx'):
                ppt = Presentation(file)
                pdf_output = BytesIO()

                # Create a ReportLab canvas for the PDF
                c = canvas.Canvas(pdf_output)
                
                # Convert each slide to an image and add it to the PDF
                for slide in ppt.slides:
                    slide_image = self.convert_slide_to_image(slide)
                    
                    # Save the slide image to a temporary file
                    with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmpfile:
                        slide_image.save(tmpfile, format='PNG')
                        tmpfile_path = tmpfile.name

                    # Add the image from the temporary file to the PDF
                    c.drawImage(tmpfile_path, 0, 0, width=600, height=400)  # Adjust size as needed

                    # Finish the page and move to the next slide
                    c.showPage()

                    # Delete the temporary file
                    os.remove(tmpfile_path)
                
                # Save the final PDF
                c.save()
                pdf_output.seek(0)
                
                self.error_manager.show_message(2015)

                return send_file(pdf_output, as_attachment=True, download_name='converted.pdf', mimetype='application/pdf')

        except Exception as e:
            return jsonify({"error": str(e)}), 500
        
app = ScholarKM()

if __name__ == '__main__':
  app.run(debug=True, threaded=True)

